#!/usr/bin/env python3
"""Erzeugt die AG-Kurzfassung als PowerPoint (.pptx).
Nutzung:  python3 scripts/make-ag-pptx.py
Benötigt: python-pptx  (pip3 install python-pptx)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# --- Vereinsfarben ---
BLUE = RGBColor(0x00, 0x2F, 0x65)
BLUE_DARK = RGBColor(0x00, 0x19, 0x3A)
YELLOW = RGBColor(0xFF, 0xDE, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1C, 0x25, 0x30)
MUTED = RGBColor(0x5B, 0x66, 0x75)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def set_run(r, text, size, color, bold=False):
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Calibri"


def accent_bar(slide):
    # gelber Balken oben
    bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.18))
    bar.fill.solid(); bar.fill.fore_color.rgb = YELLOW
    bar.line.fill.background()


def content_slide(title, bullets, kicker="MERSCH75 · Vereins-OS"):
    s = prs.slides.add_slide(BLANK)
    add_bg(s, WHITE)
    accent_bar(s)
    # Kicker
    _, ktf = textbox(s, Inches(0.7), Inches(0.45), Inches(12), Inches(0.4))
    set_run(ktf.paragraphs[0].add_run(), kicker.upper(), 12, MUTED, True)
    # Titel
    _, ttf = textbox(s, Inches(0.7), Inches(0.8), Inches(12), Inches(1.0))
    set_run(ttf.paragraphs[0].add_run(), title, 30, BLUE, True)
    # Bullets
    _, btf = textbox(s, Inches(0.8), Inches(1.95), Inches(11.8), Inches(5.0))
    for i, (txt, lvl) in enumerate(bullets):
        p = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        p.level = lvl
        p.space_after = Pt(10)
        bullet = "•  " if lvl == 0 else "–  "
        run = p.add_run()
        set_run(run, bullet + txt, 20 if lvl == 0 else 17,
                INK if lvl == 0 else MUTED, lvl == 0)
    return s


# ---------- Folie 1: Titel ----------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLUE)
band = s.shapes.add_shape(1, 0, Inches(3.05), SW, Inches(0.12))
band.fill.solid(); band.fill.fore_color.rgb = YELLOW; band.line.fill.background()
_, tf = textbox(s, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.2))
set_run(tf.paragraphs[0].add_run(), "Vereins-OS / M75-Manager", 44, WHITE, True)
_, tf2 = textbox(s, Inches(0.9), Inches(3.35), Inches(11.5), Inches(0.8))
set_run(tf2.paragraphs[0].add_run(),
        "Das eigene digitale Betriebssystem des MERSCH75", 22, YELLOW, False)
_, tf3 = textbox(s, Inches(0.9), Inches(6.2), Inches(11.5), Inches(0.6))
set_run(tf3.paragraphs[0].add_run(),
        "Generalversammlung · Kurzfassung · 29.06.2026", 14,
        RGBColor(0xC9, 0xD4, 0xE6), False)

# ---------- Folie 2: Was ist das? ----------
content_slide("Was ist das Vereins-OS?", [
    ("Ein vereinseigenes System für die gesamte Verwaltung – an einem Ort.", 0),
    ("Mitglieder, Finanzen, Sport, Kommunikation, Organisation.", 1),
    ("Selbst gehostet: alle Daten gehören dem Verein (RGPD-konform).", 0),
    ("Unabhängig von teuren Fremdanbietern.", 1),
    ("Schluss mit verteilten Excel-Listen – eine verlässliche Datenquelle.", 0),
])

# ---------- Folie 3: Nutzen ----------
content_slide("Nutzen für den Verein", [
    ("Digitale Souveränität – Daten bleiben beim Verein.", 0),
    ("Zeitersparnis durch Automatisierung (Beiträge, Anwesenheit, Spielpläne).", 0),
    ("Geringe Kosten – Open Source statt Lizenzgebühren pro Mitglied.", 0),
    ("Übergabefähig – auch ein anderes Komiteemitglied kann es pflegen.", 0),
])

# ---------- Folie 4: Stand ----------
content_slide("Wo stehen wir?", [
    ("FERTIG & im Einsatz:", 0),
    ("Sport, Mitglieder, Beiträge, Check-in, Dokumente, Kalender, E-Mail.", 1),
    ("IN ARBEIT:", 0),
    ("Chat, Profile, Website-Anbindung, Sponsoren, Galerie.", 1),
    ("GEPLANT:", 0),
    ("KI-Hilfe, Monitoring, automatisches Backup, Messenger (Matrix).", 1),
    ("Das teure Fundament steht – Erweiterungen folgen Schritt für Schritt.", 0),
])

# ---------- Folie 5: Zahlen ----------
content_slide("Größenordnung & Kosten", [
    ("Großes System: ~16.700 Code-Zeilen, 80 Datenbereiche, 233 Funktionen.", 0),
    ("Hosting ca. 60–250 € / Jahr · Software-Lizenzen 0 € (Open Source).", 0),
    ("Vergleich: kommerzielle Vereins-Software oft mehrere hundert €/Jahr.", 1),
    ("Hauptaufwand ist ehrenamtliche Zeit, nicht Geld.", 0),
])

# ---------- Folie 6: Risiken ----------
content_slide("Risiken – und wie wir sie beherrschen", [
    ("Abhängigkeit von einer Person → 2. Person einarbeiten.", 0),
    ("Sicherheit (Demo-Passwörter) → vor Echtbetrieb bereinigen.", 0),
    ("Noch kein Auto-Backup → Backup-Konzept zuerst umsetzen.", 0),
    ("Manche Module noch leer -> klar als 'in Entwicklung' kennzeichnen.", 0),
])

# ---------- Folie 7: Beschluss ----------
content_slide("Antrag an die Versammlung", [
    ("Grundsatz: Verein setzt auf das eigene Vereins-OS.", 0),
    ("Budgetfreigabe für Server-Hosting (kleiner Jahresbetrag).", 0),
    ("Start Pilotbetrieb: „Mitglieder & Beiträge\".", 0),
    ("Zweite Person zur Einarbeitung benennen (Ausfallsicherheit).", 0),
])

# ---------- Folie 8: Abschluss ----------
s = prs.slides.add_slide(BLANK)
add_bg(s, BLUE)
band = s.shapes.add_shape(1, 0, Inches(4.4), SW, Inches(0.12))
band.fill.solid(); band.fill.fore_color.rgb = YELLOW; band.line.fill.background()
_, tf = textbox(s, Inches(0.9), Inches(2.7), Inches(11.5), Inches(1.5))
set_run(tf.paragraphs[0].add_run(),
        "Ein eigenes System. Volle Datenhoheit.", 34, WHITE, True)
p = tf.paragraphs[0]
_, tf2 = textbox(s, Inches(0.9), Inches(4.8), Inches(11.5), Inches(0.8))
set_run(tf2.paragraphs[0].add_run(), "Zesumme Staark – MERSCH75", 22, YELLOW, False)

prs.save("M75-Manager-AG-Kurzfassung.pptx")
print("PPTX gespeichert: M75-Manager-AG-Kurzfassung.pptx (", len(prs.slides._sldIdLst), "Folien )")
